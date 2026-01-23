#!/usr/bin/env python3
"""
Generate PDF and PPTX from inlier visualization PNGs.
Each page/slide contains one visualization, sized for standard PowerPoint slides (16:9).

Usage:
    python generate_inliers_pdf.py --vis-dir /path/to/visualizations
    python generate_inliers_pdf.py  # Uses default paths, generates both PDF and PPTX
    python generate_inliers_pdf.py --quality 60 --max-width 1920  # Compressed
    python generate_inliers_pdf.py --format pdf  # PDF only
    python generate_inliers_pdf.py --format pptx  # PPTX only
"""

import os
import io
import argparse
import tempfile
from PIL import Image
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from natsort import natsorted


# Standard PowerPoint slide sizes (in inches)
PPTX_WIDESCREEN_16_9 = (13.333 * inch, 7.5 * inch)  # 16:9 widescreen for reportlab
PPTX_STANDARD_4_3 = (10 * inch, 7.5 * inch)  # 4:3 standard for reportlab

# PowerPoint slide dimensions in inches
PPTX_WIDTH_16_9 = 13.333
PPTX_HEIGHT_16_9 = 7.5
PPTX_WIDTH_4_3 = 10.0
PPTX_HEIGHT_4_3 = 7.5


def get_png_files(vis_dir):
    """Get sorted list of PNG files in directory."""
    if not os.path.exists(vis_dir):
        raise FileNotFoundError(f"Directory not found: {vis_dir}")

    png_files = [f for f in os.listdir(vis_dir) if f.endswith('.png')]
    return natsorted(png_files)


def compress_image(img, max_width=1920, quality=70):
    """
    Compress image by resizing and converting to JPEG.

    Args:
        img: PIL Image
        max_width: Maximum width in pixels (maintains aspect ratio)
        quality: JPEG quality (1-100, lower = smaller file)

    Returns:
        BytesIO buffer containing compressed JPEG
    """
    # Resize if larger than max_width
    if img.width > max_width:
        ratio = max_width / img.width
        new_size = (max_width, int(img.height * ratio))
        img = img.resize(new_size, Image.LANCZOS)

    # Convert to RGB if necessary (for JPEG)
    if img.mode in ('RGBA', 'P'):
        img = img.convert('RGB')

    # Save to buffer as JPEG
    buffer = io.BytesIO()
    img.save(buffer, format='JPEG', quality=quality, optimize=True)
    buffer.seek(0)

    return buffer, img.size


def generate_pdf(vis_dir, output_path, page_size=PPTX_WIDESCREEN_16_9, margin=0.25*inch,
                 max_width=None, quality=85):
    """
    Generate PDF with one visualization per page.

    Args:
        vis_dir: Directory containing PNG visualizations
        output_path: Output PDF path
        page_size: Page size tuple (width, height)
        margin: Margin around images
        max_width: Max image width for compression (None = no resize)
        quality: JPEG quality for compression (1-100)
    """
    png_files = get_png_files(vis_dir)

    if not png_files:
        print(f"No PNG files found in {vis_dir}")
        return

    compress = max_width is not None or quality < 100

    print(f"Found {len(png_files)} PNG files")
    print(f"Page size: {page_size[0]/inch:.2f}\" x {page_size[1]/inch:.2f}\" (16:9 widescreen)")
    if compress:
        print(f"Compression: max_width={max_width}px, quality={quality}")
    print(f"Output: {output_path}")

    # Create PDF
    c = canvas.Canvas(output_path, pagesize=page_size)
    page_width, page_height = page_size

    # Available space for image (with margins)
    available_width = page_width - 2 * margin
    available_height = page_height - 2 * margin

    # Use temp directory for compressed images
    with tempfile.TemporaryDirectory() as tmp_dir:
        for i, png_file in enumerate(png_files):
            png_path = os.path.join(vis_dir, png_file)

            # Extract scene_id and frame_idx from filename
            parts = png_file.replace('.png', '').split('_')
            scene_id = parts[0]
            frame_idx = '_'.join(parts[1:]) if len(parts) > 1 else ""

            # Load and optionally compress image
            with Image.open(png_path) as img:
                if compress:
                    buffer, (img_width, img_height) = compress_image(
                        img, max_width=max_width or img.width, quality=quality
                    )
                    # Save to temp file (reportlab needs file path)
                    tmp_path = os.path.join(tmp_dir, f"tmp_{i}.jpg")
                    with open(tmp_path, 'wb') as f:
                        f.write(buffer.getvalue())
                    image_path = tmp_path
                else:
                    img_width, img_height = img.size
                    image_path = png_path

            # Calculate scaling to fit in available space while maintaining aspect ratio
            scale_w = available_width / img_width
            scale_h = available_height / img_height
            scale = min(scale_w, scale_h)

            # Final image dimensions
            final_width = img_width * scale
            final_height = img_height * scale

            # Center image on page
            x = (page_width - final_width) / 2
            y = (page_height - final_height) / 2

            # Draw image
            c.drawImage(image_path, x, y, width=final_width, height=final_height)

            # Add page number at bottom
            c.setFont("Helvetica", 8)
            c.drawCentredString(page_width / 2, margin / 2,
                               f"Page {i+1}/{len(png_files)} | {scene_id} / {frame_idx}")

            # New page (except for last)
            if i < len(png_files) - 1:
                c.showPage()

            if (i + 1) % 10 == 0:
                print(f"  Processed {i+1}/{len(png_files)} pages...")

        c.save()

    # Report file size
    file_size = os.path.getsize(output_path)
    size_mb = file_size / (1024 * 1024)
    print(f"\nSaved PDF to {output_path}")
    print(f"Total pages: {len(png_files)}")
    print(f"File size: {size_mb:.1f} MB")


def generate_pptx(vis_dir, output_path, aspect="16:9", max_width=None, quality=85):
    """
    Generate PPTX with one visualization per slide.

    Args:
        vis_dir: Directory containing PNG visualizations
        output_path: Output PPTX path
        aspect: Aspect ratio ("16:9" or "4:3")
        max_width: Max image width for compression (None = no resize)
        quality: JPEG quality for compression (1-100)
    """
    png_files = get_png_files(vis_dir)

    if not png_files:
        print(f"No PNG files found in {vis_dir}")
        return

    compress = max_width is not None or quality < 100

    # Set slide dimensions
    if aspect == "16:9":
        slide_width, slide_height = PPTX_WIDTH_16_9, PPTX_HEIGHT_16_9
    else:
        slide_width, slide_height = PPTX_WIDTH_4_3, PPTX_HEIGHT_4_3

    print(f"Found {len(png_files)} PNG files")
    print(f"Slide size: {slide_width:.2f}\" x {slide_height:.2f}\" ({aspect})")
    if compress:
        print(f"Compression: max_width={max_width}px, quality={quality}")
    print(f"Output: {output_path}")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank slide

    margin = 0.25  # inches

    # Use temp directory for compressed images
    with tempfile.TemporaryDirectory() as tmp_dir:
        for i, png_file in enumerate(png_files):
            png_path = os.path.join(vis_dir, png_file)

            # Extract scene_id and frame_idx from filename
            parts = png_file.replace('.png', '').split('_')
            scene_id = parts[0]
            frame_idx = '_'.join(parts[1:]) if len(parts) > 1 else ""

            # Load and optionally compress image
            with Image.open(png_path) as img:
                if compress:
                    buffer, (img_width, img_height) = compress_image(
                        img, max_width=max_width or img.width, quality=quality
                    )
                    tmp_path = os.path.join(tmp_dir, f"tmp_{i}.jpg")
                    with open(tmp_path, 'wb') as f:
                        f.write(buffer.getvalue())
                    image_path = tmp_path
                else:
                    img_width, img_height = img.size
                    image_path = png_path

            # Add slide
            slide = prs.slides.add_slide(blank_layout)

            # Calculate image size to fit slide while maintaining aspect ratio
            available_width = slide_width - 2 * margin
            available_height = slide_height - 2 * margin - 0.3  # Extra space for caption

            img_aspect = img_width / img_height
            available_aspect = available_width / available_height

            if img_aspect > available_aspect:
                # Image is wider - fit to width
                final_width = available_width
                final_height = available_width / img_aspect
            else:
                # Image is taller - fit to height
                final_height = available_height
                final_width = available_height * img_aspect

            # Center image
            left = (slide_width - final_width) / 2
            top = (slide_height - final_height - 0.3) / 2

            # Add image
            slide.shapes.add_picture(
                image_path,
                Inches(left), Inches(top),
                Inches(final_width), Inches(final_height)
            )

            # Add caption at bottom
            caption_box = slide.shapes.add_textbox(
                Inches(0), Inches(slide_height - 0.35),
                Inches(slide_width), Inches(0.3)
            )
            tf = caption_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Page {i+1}/{len(png_files)} | {scene_id} / {frame_idx}"
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

            if (i + 1) % 10 == 0:
                print(f"  Processed {i+1}/{len(png_files)} slides...")

    prs.save(output_path)

    # Report file size
    file_size = os.path.getsize(output_path)
    size_mb = file_size / (1024 * 1024)
    print(f"\nSaved PPTX to {output_path}")
    print(f"Total slides: {len(png_files)}")
    print(f"File size: {size_mb:.1f} MB")


def detect_threshold_subdirs(base_dir):
    """Detect threshold subdirectories (e.g., 0.1cm, 0.5cm, 1.0cm)."""
    if not os.path.exists(base_dir):
        return []

    subdirs = []
    for name in os.listdir(base_dir):
        path = os.path.join(base_dir, name)
        if os.path.isdir(path) and name.endswith('cm'):
            # Check if it contains PNG files
            png_files = [f for f in os.listdir(path) if f.endswith('.png')]
            if png_files:
                subdirs.append(name)

    return natsorted(subdirs)


def generate_all(base_dir, page_size, aspect, max_width, quality, formats):
    """Generate PDF and/or PPTX for all threshold subdirectories."""
    threshold_dirs = detect_threshold_subdirs(base_dir)

    if not threshold_dirs:
        # No threshold subdirs - check if PNGs are directly in base_dir
        png_files = [f for f in os.listdir(base_dir) if f.endswith('.png')]
        if png_files:
            print(f"No threshold subdirectories found. Processing base directory...")
            dir_name = os.path.basename(base_dir.rstrip('/'))
            if 'pdf' in formats:
                output_path = os.path.join(base_dir, f"inliers_{dir_name}.pdf")
                generate_pdf(base_dir, output_path, page_size=page_size,
                            max_width=max_width, quality=quality)
            if 'pptx' in formats:
                output_path = os.path.join(base_dir, f"inliers_{dir_name}.pptx")
                generate_pptx(base_dir, output_path, aspect=aspect,
                             max_width=max_width, quality=quality)
        else:
            print(f"No PNG files found in {base_dir}")
        return

    print(f"Found {len(threshold_dirs)} threshold directories: {threshold_dirs}")
    print(f"Generating formats: {formats}")
    print("=" * 60)

    generated_files = []

    for threshold_dir in threshold_dirs:
        vis_dir = os.path.join(base_dir, threshold_dir)

        print(f"\n>>> Processing threshold: {threshold_dir}")

        if 'pdf' in formats:
            output_path = os.path.join(base_dir, f"inliers_{threshold_dir}.pdf")
            generate_pdf(vis_dir, output_path, page_size=page_size,
                        max_width=max_width, quality=quality)
            generated_files.append(output_path)

        if 'pptx' in formats:
            output_path = os.path.join(base_dir, f"inliers_{threshold_dir}.pptx")
            generate_pptx(vis_dir, output_path, aspect=aspect,
                         max_width=max_width, quality=quality)
            generated_files.append(output_path)

    print("\n" + "=" * 60)
    print(f"Generated {len(generated_files)} files in {base_dir}")
    for filepath in sorted(generated_files):
        if os.path.exists(filepath):
            size_mb = os.path.getsize(filepath) / (1024 * 1024)
            print(f"  - {os.path.basename(filepath)} ({size_mb:.1f} MB)")


def main():
    parser = argparse.ArgumentParser(description="Generate PDF and/or PPTX from inlier visualization PNGs")
    parser.add_argument("--vis-dir", type=str,
                       default="/cluster/scratch/aoezkan/planeseg/scannetpp/visualizations/inliers/baselines_n20_seed42_v5",
                       help="Base directory containing threshold subdirectories (or PNGs directly)")
    parser.add_argument("--output", type=str, default=None,
                       help="Output path (only used with --threshold, extension determines format)")
    parser.add_argument("--aspect", type=str, choices=["16:9", "4:3"], default="16:9",
                       help="Page/slide aspect ratio (default: 16:9 widescreen)")
    parser.add_argument("--threshold", type=str, default=None,
                       help="Process only this threshold (e.g., '0.5cm'). If not specified, process all.")
    parser.add_argument("--format", type=str, choices=["pdf", "pptx", "both"], default="both",
                       help="Output format (default: both)")
    parser.add_argument("--quality", type=int, default=70,
                       help="JPEG quality for compression (1-100, default: 70). Lower = smaller file.")
    parser.add_argument("--max-width", type=int, default=1920,
                       help="Max image width in pixels (default: 1920). Set to 0 to disable resizing.")
    parser.add_argument("--no-compress", action="store_true",
                       help="Disable compression (use original PNGs)")

    args = parser.parse_args()

    # Select page size (for PDF)
    page_size = PPTX_WIDESCREEN_16_9 if args.aspect == "16:9" else PPTX_STANDARD_4_3

    # Compression settings
    if args.no_compress:
        max_width = None
        quality = 100
    else:
        max_width = args.max_width if args.max_width > 0 else None
        quality = args.quality

    # Determine formats to generate
    if args.format == "both":
        formats = ['pdf', 'pptx']
    else:
        formats = [args.format]

    # Single threshold or all thresholds
    if args.threshold:
        # Process single threshold
        vis_dir = os.path.join(args.vis_dir, args.threshold)

        if args.output:
            # Use specified output path
            if 'pdf' in formats:
                output_pdf = args.output if args.output.endswith('.pdf') else args.output + '.pdf'
                generate_pdf(vis_dir, output_pdf, page_size=page_size,
                            max_width=max_width, quality=quality)
            if 'pptx' in formats:
                output_pptx = args.output if args.output.endswith('.pptx') else args.output.replace('.pdf', '') + '.pptx'
                generate_pptx(vis_dir, output_pptx, aspect=args.aspect,
                             max_width=max_width, quality=quality)
        else:
            # Default output paths
            if 'pdf' in formats:
                output_path = os.path.join(args.vis_dir, f"inliers_{args.threshold}.pdf")
                generate_pdf(vis_dir, output_path, page_size=page_size,
                            max_width=max_width, quality=quality)
            if 'pptx' in formats:
                output_path = os.path.join(args.vis_dir, f"inliers_{args.threshold}.pptx")
                generate_pptx(vis_dir, output_path, aspect=args.aspect,
                             max_width=max_width, quality=quality)
    else:
        # Auto-detect and process all thresholds
        generate_all(args.vis_dir, page_size=page_size, aspect=args.aspect,
                    max_width=max_width, quality=quality, formats=formats)


if __name__ == "__main__":
    main()
